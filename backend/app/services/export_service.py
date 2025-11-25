from io import BytesIO
from docx import Document
from pptx import Presentation
from pptx.util import Inches

def assemble_docx(project, nodes):
    doc = Document()
    doc.add_heading(project.title or "Document", level=1)
    if project.main_prompt:
        doc.add_paragraph(f"Topic: {project.main_prompt}")
    for n in sorted(nodes, key=lambda x: x.idx):
        doc.add_heading(n.title or "", level=2)
        if n.content_current:
            for para in (n.content_current.split("\n")):
                if para.strip().startswith("- ") or para.strip().startswith("* "):
                    p = doc.add_paragraph(para.strip().lstrip("-* ").strip(), style='List Bullet')
                else:
                    doc.add_paragraph(para)
    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

def assemble_pptx(project, nodes):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    if project.title:
        slide.shapes.title.text = project.title
    if project.main_prompt:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = project.main_prompt
    for n in sorted(nodes, key=lambda x: x.idx):
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        title = s.shapes.title
        body = s.shapes.placeholders[1].text_frame
        title.text = n.title or ""
        if n.content_current:
            for li, para in enumerate(n.content_current.split("\n")):
                if li==0:
                    body.text = para
                else:
                    body.add_paragraph().text = para
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio
